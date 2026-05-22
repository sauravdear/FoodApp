from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ── Palette ─────────────────────────────────────────────────────────────────
BG        = RGBColor(0x0A, 0x0F, 0x1E)   # deep navy
BG2       = RGBColor(0x0D, 0x14, 0x2B)
PANEL     = RGBColor(0x11, 0x1C, 0x36)
CARD      = RGBColor(0x16, 0x23, 0x3A)
BORDER    = RGBColor(0x1E, 0x30, 0x50)
ACCENT    = RGBColor(0x00, 0xD4, 0x8A)   # teal-green
ACCENT2   = RGBColor(0x00, 0xA8, 0xFF)   # electric blue
GOLD      = RGBColor(0xF5, 0xA6, 0x23)
RED       = RGBColor(0xFF, 0x4C, 0x4C)
PURPLE    = RGBColor(0xA7, 0x8B, 0xFA)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE  = RGBColor(0xE2, 0xE8, 0xF0)
GRAY      = RGBColor(0x64, 0x74, 0x8B)
LGRAY     = RGBColor(0x94, 0xA3, 0xB8)

TEAM_NAME    = "Team  6"
TEAM_MEMBERS = "Saurav  |  Tushar  |  Ankit  |  Samyank"
PROJECT      = "FoodRedist"

# ── Low-level helpers ────────────────────────────────────────────────────────
def rect(slide, l, t, w, h, fill, alpha=None, line_clr=None, line_w=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line_clr:
        s.line.color.rgb = line_clr
        if line_w: s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=16, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color
    return tb

def add_p(tf, text, size=14, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, before=4, italic=False, font=None):
    p = tf.add_paragraph(); p.alignment = align; p.space_before = Pt(before)
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color
    if font: run.font.name = font
    return p

def bg(slide):
    rect(slide, 0, 0, 13.33, 7.5, BG)

def footer(slide, n, total=14):
    rect(slide, 0, 7.22, 13.33, 0.28, PANEL)
    rect(slide, 0, 7.22, 13.33, 0.025, ACCENT)
    txt(slide, TEAM_NAME, 0.3, 7.25, 3, 0.22, size=10, color=LGRAY)
    txt(slide, PROJECT + " — Perishable Food Stock Redistribution System",
        3, 7.25, 7.3, 0.22, size=10, color=GRAY, align=PP_ALIGN.CENTER)
    txt(slide, f"{n} / {total}", 12.3, 7.25, 0.8, 0.22,
        size=10, color=LGRAY, align=PP_ALIGN.RIGHT)

def section_header(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 0.06, ACCENT)
    txt(slide, title, 0.45, 0.18, 11, 0.72,
        size=32, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.45, 0.88, 12, 0.42,
            size=15, color=LGRAY, italic=True)
    rect(slide, 0.45, 1.35, 1.2, 0.05, ACCENT)

def bullet_box(slide, l, t, w, h, heading, items, accent=ACCENT,
               head_size=15, item_size=13.5, icon="▸"):
    rect(slide, l, t, w, h, CARD, line_clr=BORDER, line_w=0.5)
    rect(slide, l, t, 0.05, h, accent)
    txt(slide, heading, l+0.2, t+0.15, w-0.3, 0.38,
        size=head_size, bold=True, color=accent)
    tb = slide.shapes.add_textbox(
        Inches(l+0.2), Inches(t+0.58), Inches(w-0.3), Inches(h-0.72))
    tb.word_wrap = True; tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run(); run.text = f"{icon}  {item}"
        run.font.size = Pt(item_size); run.font.color.rgb = OFFWHITE


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)

# left dark panel
rect(s, 0, 0, 5.8, 7.5, PANEL)
rect(s, 5.75, 0, 0.06, 7.5, ACCENT)

# decorative rings (nested rects as rings)
for i, (clr, op) in enumerate([(ACCENT, 0.08), (ACCENT2, 0.05), (ACCENT, 0.03)]):
    sz = 2.2 + i*0.7
    rect(s, (5.8-sz)/2, 1.0-i*0.35, sz, sz, PANEL,
         line_clr=clr, line_w=0.8)

txt(s, "🥗", 1.5, 1.3, 2.8, 1.5, size=72, align=PP_ALIGN.CENTER)

txt(s, PROJECT, 0.25, 3.0, 5.3, 0.9,
    size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Perishable Food Stock", 0.25, 3.88, 5.3, 0.48,
    size=18, color=OFFWHITE, align=PP_ALIGN.CENTER)
txt(s, "Redistribution System", 0.25, 4.32, 5.3, 0.48,
    size=18, color=OFFWHITE, align=PP_ALIGN.CENTER)

rect(s, 0.5, 5.0, 4.8, 0.03, ACCENT)

txt(s, TEAM_NAME, 0.25, 5.1, 5.3, 0.42,
    size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
txt(s, TEAM_MEMBERS, 0.25, 5.52, 5.3, 0.36,
    size=13, color=LGRAY, align=PP_ALIGN.CENTER)
txt(s, "B.Tech Computer Science  |  2025–26", 0.25, 5.9, 5.3, 0.35,
    size=12, color=GRAY, align=PP_ALIGN.CENTER)

# right side
txt(s, "A Smart Approach to", 6.1, 1.4, 6.9, 0.55,
    size=22, color=LGRAY, bold=False)
txt(s, "Zero Food Waste", 6.1, 1.9, 6.9, 0.7,
    size=36, bold=True, color=ACCENT)
txt(s, "Using MERN Stack Technology", 6.1, 2.62, 6.9, 0.45,
    size=17, color=OFFWHITE)

txt(s,
    "Real-time inventory tracking, automated velocity\n"
    "analysis, and intelligent stock redistribution\n"
    "across multiple stores — before items expire.",
    6.1, 3.3, 6.8, 1.5, size=14.5, color=LGRAY)

# stat chips
for i, (val, lbl, clr) in enumerate([
    ("₹92K Cr", "Annual Food Waste India", RED),
    ("30–40%", "Perishables Wasted", GOLD),
    ("MERN", "Tech Stack Used", ACCENT),
]):
    lx = 6.1 + i * 2.38
    rect(s, lx, 5.05, 2.18, 1.25, CARD, line_clr=clr, line_w=0.8)
    txt(s, val, lx, 5.18, 2.18, 0.52,
        size=19, bold=True, color=clr, align=PP_ALIGN.CENTER)
    txt(s, lbl, lx, 5.72, 2.18, 0.48,
        size=11, color=LGRAY, align=PP_ALIGN.CENTER)

rect(s, 0, 7.22, 13.33, 0.28, PANEL)
rect(s, 0, 7.22, 13.33, 0.025, ACCENT)
txt(s, TEAM_NAME + "  |  " + TEAM_MEMBERS, 0, 7.26, 13.33, 0.22,
    size=10, color=GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s); section_header(s, "Agenda", "Presentation roadmap — what we will cover today")

agenda = [
    ("01", "Introduction",       "Overview of the project and its domain"),
    ("02", "Problem Statement",  "The core challenges in perishable food management"),
    ("03", "Existing System",    "Current approaches and their limitations"),
    ("04", "Proposed Solution",  "FoodRedist — our intelligent redistribution platform"),
    ("05", "Architecture",       "System design and component interaction"),
    ("06", "Technologies Used",  "MERN stack and supporting libraries"),
    ("07", "Implementation",     "Key modules, features and code structure"),
    ("08", "Results",            "Output, recommendations and transfer demo"),
    ("09", "Challenges",         "Technical and logistical hurdles we overcame"),
    ("10", "Future Scope",       "Enhancements planned for next phase"),
    ("11", "Conclusion",         "Key takeaways and project impact"),
    ("12", "Q & A",              "Open floor for questions"),
]

cols = [agenda[:6], agenda[6:]]
for ci, col_items in enumerate(cols):
    lx = 0.45 + ci * 6.45
    for ri, (num, title, desc) in enumerate(col_items):
        ty = 1.5 + ri * 0.92
        rect(s, lx, ty, 6.1, 0.8, CARD, line_clr=BORDER, line_w=0.4)
        rect(s, lx, ty, 0.05, 0.8, ACCENT if ri % 2 == 0 else ACCENT2)
        rect(s, lx+0.2, ty+0.1, 0.52, 0.52, PANEL)
        txt(s, num, lx+0.2, ty+0.1, 0.52, 0.52,
            size=12, bold=True, color=ACCENT if ri % 2 == 0 else ACCENT2,
            align=PP_ALIGN.CENTER)
        txt(s, title, lx+0.88, ty+0.1, 4.2, 0.32,
            size=14, bold=True, color=WHITE)
        txt(s, desc, lx+0.88, ty+0.43, 5.0, 0.28,
            size=11, color=LGRAY, italic=True)

footer(s, 2)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — INTRODUCTION
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s); section_header(s, "Introduction", "Understanding the domain — perishable food supply chain")

# Overview paragraph
rect(s, 0.45, 1.48, 12.43, 0.95, PANEL, line_clr=BORDER, line_w=0.4)
rect(s, 0.45, 1.48, 0.06, 0.95, ACCENT2)
tb = s.shapes.add_textbox(Inches(0.7), Inches(1.55), Inches(12.0), Inches(0.82))
tb.word_wrap = True; tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = ("FoodRedist is a full-stack MERN application designed to solve one of the most persistent challenges "
            "in the FMCG and grocery retail industry — perishable food items expiring unsold in slow-moving stores "
            "while fast-moving stores run out of stock. The platform tracks inventory in real-time, computes sales "
            "velocity per store, and generates intelligent transfer recommendations before items hit their expiry window.")
run.font.size = Pt(14); run.font.color.rgb = OFFWHITE

# 4 intro cards
intro_cards = [
    ("🎯", "Project Domain", ACCENT,
     ["Grocery & Retail Supply Chain", "Multi-store Inventory Management",
      "Perishable Goods Lifecycle", "Last-mile food redistribution"]),
    ("🏗️", "Core Concept", ACCENT2,
     ["Track: Real-time stock monitoring", "Analyze: Sales velocity per store",
      "Recommend: Smart transfer alerts", "Execute: 1-click stock movement"]),
    ("👥", "Target Users", GOLD,
     ["Grocery chain managers", "Warehouse supervisors",
      "Store inventory officers", "Supply chain analysts"]),
    ("📐", "Project Scope", PURPLE,
     ["Backend REST API (Node/Express)", "Mongoose data models",
      "Velocity Engine algorithm", "React dashboard with dark UI"]),
]

for i, (icon, head, clr, pts) in enumerate(intro_cards):
    lx = 0.45 + i * 3.23
    ty = 2.6
    rect(s, lx, ty, 3.05, 3.9, CARD, line_clr=BORDER, line_w=0.4)
    rect(s, lx, ty, 3.05, 0.06, clr)
    txt(s, icon, lx+0.15, ty+0.12, 0.65, 0.65, size=26)
    txt(s, head, lx+0.82, ty+0.15, 2.1, 0.42, size=14, bold=True, color=clr)
    tb2 = s.shapes.add_textbox(Inches(lx+0.18), Inches(ty+0.72), Inches(2.72), Inches(3.0))
    tb2.word_wrap = True; tf2 = tb2.text_frame; tf2.word_wrap = True
    for j, pt in enumerate(pts):
        p2 = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
        p2.space_before = Pt(6)
        run2 = p2.add_run(); run2.text = "▸  " + pt
        run2.font.size = Pt(13); run2.font.color.rgb = OFFWHITE

footer(s, 3)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — PROBLEM STATEMENT
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s); section_header(s, "Problem Statement", "Why perishable food management is broken today")

# Key stat banner
rect(s, 0.45, 1.48, 12.43, 0.82, RGBColor(0x1A, 0x08, 0x08), line_clr=RED, line_w=0.6)
txt(s, "⚠️   India wastes ₹92,000 Crore worth of food annually — 40% of all perishables never reach the consumer",
    0.65, 1.58, 12.1, 0.55, size=14.5, bold=True, color=RGBColor(0xFC, 0xA5, 0xA5))

problems = [
    ("P1", "Uneven Stock Distribution", RED,
     ["Slow stores hold excess stock that will expire",
      "Fast stores run dry causing lost sales daily",
      "No automated cross-store balancing exists"]),
    ("P2", "No Real-Time Visibility", GOLD,
     ["Inventory tracked manually via spreadsheets",
      "Expiry dates not monitored proactively",
      "Managers react after waste — not before"]),
    ("P3", "Zero Sales Velocity Insight", ACCENT2,
     ["Stores don't know their own sell-through rate",
      "Cannot predict which items will expire unsold",
      "No data to guide transfer decisions"]),
    ("P4", "Manual & Error-Prone Transfers", PURPLE,
     ["Transfer decisions based on intuition only",
      "No audit trail of stock movements",
      "Delays lead to spoilage and revenue loss"]),
]

for i, (code, head, clr, pts) in enumerate(problems):
    col = i % 2; row = i // 2
    lx = 0.45 + col * 6.45
    ty = 2.5 + row * 2.3
    rect(s, lx, ty, 6.1, 2.1, CARD, line_clr=BORDER, line_w=0.4)
    rect(s, lx, ty, 0.06, 2.1, clr)
    rect(s, lx+0.2, ty+0.15, 0.58, 0.38, clr)
    txt(s, code, lx+0.2, ty+0.15, 0.58, 0.38,
        size=11, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(s, head, lx+0.9, ty+0.15, 5.0, 0.38, size=14, bold=True, color=WHITE)
    tb = s.shapes.add_textbox(Inches(lx+0.22), Inches(ty+0.62), Inches(5.72), Inches(1.35))
    tb.word_wrap = True; tf = tb.text_frame; tf.word_wrap = True
    for j, pt in enumerate(pts):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run(); run.text = "→  " + pt
        run.font.size = Pt(13); run.font.color.rgb = LGRAY

footer(s, 4)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — EXISTING SYSTEM
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s); section_header(s, "Existing System", "Current approaches and their critical limitations")

# Comparison table header
rect(s, 0.45, 1.48, 12.43, 0.5, PANEL, line_clr=BORDER, line_w=0.4)
for lx, w, label, clr in [
    (0.45, 3.8, "EXISTING APPROACH", LGRAY),
    (4.25, 4.0, "LIMITATION", RED),
    (8.25, 4.63, "IMPACT ON BUSINESS", GOLD),
]:
    txt(s, label, lx+0.15, 1.53, w-0.2, 0.35,
        size=11, bold=True, color=clr, align=PP_ALIGN.CENTER)

rows = [
    ("Manual Spreadsheet Tracking",
     "Updated once a day — too slow for perishables",
     "Items expire before action is taken"),
    ("Visual Shelf Inspection",
     "Subjective — depends on individual attention",
     "Inconsistent; misses slow-selling batches"),
    ("Basic POS Software",
     "Records sales but gives no redistribution advice",
     "Data sits unused; no velocity analysis"),
    ("Email / Phone Coordination",
     "Slow, no audit trail, prone to miscommunication",
     "Transfers delayed by hours or days"),
    ("ERP Systems (SAP / Oracle)",
     "Expensive, rigid, not designed for perishables",
     "SMEs cannot afford; poor expiry handling"),
]

for i, (approach, limitation, impact) in enumerate(rows):
    ty = 2.08 + i * 0.98
    bg_clr = CARD if i % 2 == 0 else PANEL
    rect(s, 0.45, ty, 12.43, 0.88, bg_clr, line_clr=BORDER, line_w=0.3)
    rect(s, 0.45, ty, 0.05, 0.88, ACCENT if i % 2 == 0 else ACCENT2)
    txt(s, approach, 0.65, ty+0.12, 3.45, 0.62, size=13, bold=True, color=OFFWHITE)
    txt(s, limitation, 4.4, ty+0.12, 3.65, 0.62, size=12.5, color=RGBColor(0xFC, 0xA5, 0xA5))
    txt(s, impact, 8.3, ty+0.12, 4.35, 0.62, size=12.5, color=RGBColor(0xFD, 0xE6, 0x8A))

# Gap lines between cols
for lx in [4.25, 8.25]:
    rect(s, lx, 1.48, 0.02, 5.45, BORDER)

footer(s, 5)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PROPOSED SOLUTION
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s); section_header(s, "Proposed Solution", "FoodRedist — Intelligent, automated, real-time redistribution")

# Tagline
rect(s, 0.45, 1.48, 12.43, 0.72, RGBColor(0x06, 0x2A, 0x1F), line_clr=ACCENT, line_w=0.7)
txt(s, "✅   FoodRedist tracks inventory, computes daily sales velocity per store, flags expiring items, and "
    "recommends optimal stock transfers — all automated, all real-time.",
    0.65, 1.55, 12.05, 0.58, size=14, color=RGBColor(0xA7, 0xF3, 0xD0))

# 3 pillars
pillars = [
    ("TRACK", "Real-Time Monitoring", ACCENT,
     "📡",
     ["Every store's current stock level",
      "Expiry dates per batch & SKU",
      "Live sales velocity (units/day)",
      "Automated expiry window alerts"]),
    ("ANALYZE", "Velocity Engine", ACCENT2,
     "🧮",
     ["Formula: daysToSellOut = stock ÷ velocity",
      "Classifies each store as SLOW or FAST",
      "Computes exact surplus units per store",
      "Greedy algorithm matches slow→fast"]),
    ("EXECUTE", "Seamless Transfer", GOLD,
     "⚡",
     ["1-click approval in dashboard",
      "Atomic DB writes (no partial state)",
      "Complete TransferLog audit trail",
      "Stock counts update instantly"]),
]

for i, (tag, head, clr, icon, pts) in enumerate(pillars):
    lx = 0.45 + i * 4.3
    ty = 2.35
    rect(s, lx, ty, 4.1, 4.65, CARD, line_clr=clr, line_w=0.8)
    rect(s, lx, ty, 4.1, 0.07, clr)
    rect(s, lx+0.2, ty+0.18, 0.95, 0.42, clr)
    txt(s, tag, lx+0.2, ty+0.18, 0.95, 0.42,
        size=11, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(s, icon, lx+1.28, ty+0.12, 0.7, 0.58, size=28)
    txt(s, head, lx+0.2, ty+0.72, 3.7, 0.45, size=15, bold=True, color=clr)
    tb = s.shapes.add_textbox(Inches(lx+0.2), Inches(ty+1.28), Inches(3.72), Inches(3.22))
    tb.word_wrap = True; tf = tb.text_frame; tf.word_wrap = True
    for j, pt in enumerate(pts):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_before = Pt(7)
        run = p.add_run(); run.text = "▸  " + pt
        run.font.size = Pt(13.5); run.font.color.rgb = OFFWHITE

    if i < 2:
        txt(s, "⟶", lx+4.1, ty+2.1, 0.5, 0.5,
            size=24, color=BORDER, align=PP_ALIGN.CENTER)

footer(s, 6)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s); section_header(s, "System Architecture", "Three-tier architecture — Client, Server, Database")

# Tier diagram
tiers = [
    ("PRESENTATION TIER", "React.js Frontend", ACCENT2, "🖥️",
     ["Vite + React 19", "React Router DOM", "Tailwind CSS (dark UI)",
      "AppContext (global state)", "useFetch custom hook"]),
    ("APPLICATION TIER", "Node.js + Express", ACCENT, "⚙️",
     ["REST API endpoints", "Velocity Engine service", "Controllers & Routes",
      "CORS middleware", "Error handling layer"]),
    ("DATA TIER", "MongoDB Database", GOLD, "🗄️",
     ["FoodItem collection", "Store + Inventory", "TransferLog collection",
      "Mongoose ODM", "Atlas / Local MongoDB"]),
]

for i, (tier_lbl, tech, clr, icon, comps) in enumerate(tiers):
    lx = 0.4 + i * 4.3
    ty = 1.48
    rect(s, lx, ty, 4.1, 5.5, CARD, line_clr=clr, line_w=0.9)
    rect(s, lx, ty, 4.1, 0.07, clr)
    rect(s, lx+0.15, ty+0.12, 3.8, 0.35, PANEL)
    txt(s, tier_lbl, lx+0.15, ty+0.12, 3.8, 0.35,
        size=10, bold=True, color=clr, align=PP_ALIGN.CENTER)
    txt(s, icon, lx+1.6, ty+0.6, 0.9, 0.8, size=38, align=PP_ALIGN.CENTER)
    txt(s, tech, lx+0.15, ty+1.42, 3.8, 0.48,
        size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, lx+0.5, ty+1.98, 3.1, 0.03, clr)
    tb = s.shapes.add_textbox(Inches(lx+0.22), Inches(ty+2.1), Inches(3.66), Inches(3.22))
    tb.word_wrap = True; tf = tb.text_frame; tf.word_wrap = True
    for j, comp in enumerate(comps):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run(); run.text = "◆  " + comp
        run.font.size = Pt(13.5); run.font.color.rgb = OFFWHITE

    if i < 2:
        txt(s, "⟷", lx+4.1, ty+2.55, 0.5, 0.55,
            size=22, color=LGRAY, align=PP_ALIGN.CENTER)
        txt(s, "HTTP / REST", lx+3.92, ty+3.1, 0.88, 0.35,
            size=9, color=GRAY, align=PP_ALIGN.CENTER)

# Data flow arrow bottom
rect(s, 1.5, 7.0, 10.33, 0.02, BORDER)
txt(s, "← →  Bidirectional data flow  |  JSON over HTTP  |  Vite Proxy (dev) / CORS (prod)",
    0.45, 6.92, 12.43, 0.28,
    size=11, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

footer(s, 7)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — TECHNOLOGIES USED
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s); section_header(s, "Technologies Used", "Complete technology stack with version details")

techs = [
    ("🍃", "MongoDB", "v8.0", "Database",   ACCENT,
     "NoSQL document store. Stores FoodItem, Store & TransferLog collections. "
     "Mongoose ODM provides schema validation and query helpers."),
    ("⚡", "Express.js", "v4.19", "Backend",  GOLD,
     "Minimal Node.js web framework. Handles routing, middleware (CORS, JSON), "
     "and wires controllers to REST endpoints."),
    ("⚛️", "React.js", "v19.2", "Frontend",  ACCENT2,
     "Component-based UI library. Context API manages global state; "
     "React Router DOM v6 handles SPA navigation between Dashboard & Transfers."),
    ("🟩", "Node.js", "v20 LTS", "Runtime",  RGBColor(0x68, 0xA0, 0x63),
     "JavaScript runtime powering the Express server. ES Modules (import/export) "
     "used throughout for clean modular code."),
    ("⚡", "Vite", "v8.0", "Build Tool",    RGBColor(0xBD, 0x34, 0xFE),
     "Lightning-fast dev server with HMR. Configured with a /api proxy to "
     "forward frontend requests to the backend — no CORS issues in dev."),
    ("🎨", "Tailwind CSS", "v3.4", "Styling", RGBColor(0x38, 0xBD, 0xF8),
     "Utility-first CSS framework. Enables the dark theme dashboard with "
     "consistent spacing, colors, and responsive layout without custom CSS."),
]

for i, (icon, name, ver, role, clr, desc) in enumerate(techs):
    col = i % 2; row = i // 2
    lx = 0.45 + col * 6.45
    ty = 1.52 + row * 1.92
    rect(s, lx, ty, 6.1, 1.75, CARD, line_clr=BORDER, line_w=0.4)
    rect(s, lx, ty, 0.06, 1.75, clr)
    txt(s, icon, lx+0.18, ty+0.3, 0.72, 0.72, size=30)
    txt(s, name, lx+0.95, ty+0.14, 2.5, 0.42, size=17, bold=True, color=clr)
    rect(s, lx+3.5, ty+0.18, 0.95, 0.35, PANEL)
    txt(s, ver, lx+3.5, ty+0.18, 0.95, 0.35,
        size=11, bold=True, color=LGRAY, align=PP_ALIGN.CENTER)
    txt(s, role, lx+0.95, ty+0.55, 2.5, 0.32, size=12, color=GRAY, italic=True)
    txt(s, desc, lx+0.18, ty+0.98, 5.76, 0.72, size=12.5, color=OFFWHITE, wrap=True)

footer(s, 8)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — IMPLEMENTATION
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s); section_header(s, "Implementation", "Key modules, folder structure and core algorithm")

# Left — folder tree
rect(s, 0.45, 1.48, 4.5, 5.65, CARD, line_clr=BORDER, line_w=0.4)
rect(s, 0.45, 1.48, 0.06, 5.65, ACCENT)
txt(s, "📁  Project Structure", 0.68, 1.55, 4.1, 0.4, size=13, bold=True, color=ACCENT)

tree_lines = [
    ("food-redistribution-system/", False, WHITE),
    ("├── backend/", False, ACCENT2),
    ("│   ├── config/db.js", True, LGRAY),
    ("│   ├── models/", False, GOLD),
    ("│   │   ├── FoodItem.js", True, LGRAY),
    ("│   │   ├── Store.js", True, LGRAY),
    ("│   │   └── TransferLog.js", True, LGRAY),
    ("│   ├── controllers/", False, GOLD),
    ("│   │   └── inventoryController.js", True, LGRAY),
    ("│   ├── services/", False, ACCENT),
    ("│   │   └── velocityEngine.js ⭐", True, ACCENT),
    ("│   ├── routes/", False, GOLD),
    ("│   └── server.js", True, LGRAY),
    ("└── frontend/src/", False, ACCENT2),
    ("    ├── components/dashboard/", True, LGRAY),
    ("    ├── context/AppContext.jsx", True, LGRAY),
    ("    └── pages/", True, LGRAY),
]
tb = s.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(4.22), Inches(5.0))
tb.word_wrap = False; tf = tb.text_frame; tf.word_wrap = False
for j, (line, indent, clr) in enumerate(tree_lines):
    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
    p.space_before = Pt(3)
    run = p.add_run(); run.text = line
    run.font.size = Pt(11.5); run.font.name = "Consolas"; run.font.color.rgb = clr

# Right — core algo + API
rect(s, 5.15, 1.48, 7.73, 2.6, CARD, line_clr=BORDER, line_w=0.4)
rect(s, 5.15, 1.48, 0.06, 2.6, ACCENT2)
txt(s, "🧮  Velocity Engine Algorithm", 5.38, 1.55, 7.2, 0.4, size=13, bold=True, color=ACCENT2)

algo = [
    "1.  Fetch FoodItems expiring within N days (default: 3)",
    "2.  For each item, query all stores carrying that item",
    "3.  Compute:  daysToSellOut = currentStock / salesVelocity",
    "4.  SLOW store:  daysToSellOut > daysUntilExpiry",
    "     surplus = stock − (velocity × daysLeft)",
    "5.  FAST store:  velocity ≥ threshold AND stock ≤ guard",
    "     demandGap = (velocity × daysLeft) − stock",
    "6.  Greedy match:  slow surplus  →  fast demand gap",
]
tb2 = s.shapes.add_textbox(Inches(5.32), Inches(2.0), Inches(7.38), Inches(2.0))
tb2.word_wrap = True; tf2 = tb2.text_frame; tf2.word_wrap = True
for j, line in enumerate(algo):
    p = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
    p.space_before = Pt(4)
    run = p.add_run(); run.text = line
    run.font.size = Pt(12); run.font.name = "Consolas"
    run.font.color.rgb = LGRAY if line.startswith("     ") else OFFWHITE

# API summary
rect(s, 5.15, 4.25, 7.73, 2.88, CARD, line_clr=BORDER, line_w=0.4)
rect(s, 5.15, 4.25, 0.06, 2.88, GOLD)
txt(s, "🌐  REST API Endpoints", 5.38, 4.32, 7.2, 0.4, size=13, bold=True, color=GOLD)

apis = [
    ("GET",   "/api/inventory/recommendations", ACCENT, "Run velocity engine"),
    ("POST",  "/api/inventory/transfer",        GOLD,   "Execute stock transfer"),
    ("GET",   "/api/inventory/transfers",       ACCENT2,"Transfer history"),
    ("GET",   "/api/stores",                    PURPLE, "All stores + inventory"),
    ("PATCH", "/api/stores/:id/velocity",       LGRAY,  "Update sales velocity"),
]
for j, (method, path, clr, desc) in enumerate(apis):
    ty2 = 4.75 + j * 0.46
    rect(s, 5.32, ty2, 0.72, 0.32, clr)
    txt(s, method, 5.32, ty2, 0.72, 0.32, size=10, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(s, path, 6.1, ty2+0.01, 3.6, 0.3, size=11, color=WHITE)
    txt(s, desc, 9.72, ty2+0.01, 2.9, 0.3, size=11, color=LGRAY, italic=True)

footer(s, 9)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — RESULTS
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s); section_header(s, "Results", "System output — recommendations generated from seed data")

# Stat row
stats = [
    ("4",  "Expiring Items\nDetected",  ACCENT),
    ("3",  "Stores\nAnalyzed",          ACCENT2),
    ("7",  "Transfer\nRecommendations", GOLD),
    ("0",  "Items\nWasted",             RGBColor(0x34, 0xD3, 0x99)),
]
for i, (val, lbl, clr) in enumerate(stats):
    lx = 0.45 + i * 3.23
    rect(s, lx, 1.48, 3.05, 1.38, CARD, line_clr=clr, line_w=0.9)
    rect(s, lx, 1.48, 3.05, 0.07, clr)
    txt(s, val, lx, 1.62, 3.05, 0.72,
        size=44, bold=True, color=clr, align=PP_ALIGN.CENTER)
    txt(s, lbl, lx, 2.32, 3.05, 0.48,
        size=13, color=LGRAY, align=PP_ALIGN.CENTER)

# Sample output table
rect(s, 0.45, 3.05, 12.43, 0.42, PANEL, line_clr=BORDER, line_w=0.4)
for lx2, w2, lbl2 in [(0.45,2.8,"ITEM"),(3.25,2.8,"FROM (SLOW)"),(6.05,2.8,"TO (FAST)"),
                      (8.85,1.7,"QTY"),(10.55,2.33,"EXPIRES IN")]:
    txt(s, lbl2, lx2+0.12, 3.1, w2-0.1, 0.3,
        size=11, bold=True, color=LGRAY, align=PP_ALIGN.CENTER)

results = [
    ("Greek Yogurt",    "Downtown Fresh (1 u/d)",  "North Side Market (30 u/d)", "37",  "1.5 days", RED),
    ("Greek Yogurt",    "Downtown Fresh (1 u/d)",  "West Loop Grocer (15 u/d)",  "11",  "1.5 days", RED),
    ("Whole Milk",      "Downtown Fresh (2 u/d)",  "North Side Market (40 u/d)", "70",  "2.0 days", GOLD),
    ("Whole Milk",      "Downtown Fresh (2 u/d)",  "West Loop Grocer (18 u/d)",  "21",  "2.0 days", GOLD),
    ("Cheddar Cheese",  "Downtown Fresh (1 u/d)",  "North Side Market (20 u/d)", "45",  "2.5 days", GOLD),
    ("Sourdough Bread", "Downtown Fresh (0 u/d)",  "North Side Market (25 u/d)", "19",  "1.0 day",  RED),
    ("Sourdough Bread", "Downtown Fresh (0 u/d)",  "West Loop Grocer (22 u/d)",  "18",  "1.0 day",  RED),
]
for i, (item, frm, to, qty, exp, urg) in enumerate(results):
    ty2 = 3.57 + i * 0.49
    bg_c = CARD if i % 2 == 0 else PANEL
    rect(s, 0.45, ty2, 12.43, 0.45, bg_c, line_clr=BORDER, line_w=0.2)
    rect(s, 0.45, ty2, 0.05, 0.45, urg)
    for lx2, w2, val, clr2 in [
        (0.5, 2.75, item,  WHITE),
        (3.25, 2.75, frm,  OFFWHITE),
        (6.05, 2.75, to,   OFFWHITE),
        (8.85, 1.65, qty,  ACCENT),
        (10.55, 2.28, exp, urg),
    ]:
        txt(s, val, lx2+0.08, ty2+0.07, w2-0.1, 0.32, size=12,
            color=clr2, align=PP_ALIGN.CENTER)

footer(s, 10)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — CHALLENGES
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s); section_header(s, "Challenges Faced", "Technical and design hurdles encountered during development")

challenges = [
    ("⚙️", "MongoDB Transactions on Standalone", RED,
     "Problem", "Resolution",
     "Mongoose session transactions require a replica set. Local standalone MongoDB threw: "
     "'Transaction numbers only allowed on replica set member'",
     "Replaced session-based transaction with sequential atomic $inc operations. "
     "Each write uses its own MongoDB query-level atomicity. Atlas (replica set) will support full transactions."),
    ("🌐", "Atlas DNS / Network Block", GOLD,
     "Problem", "Resolution",
     "querySrv ECONNREFUSED on MongoDB Atlas SRV record. Port 27017 blocked by ISP/firewall — "
     "Test-NetConnection returned False.",
     "Switched to local MongoDB for development. Atlas URI preserved in .env for production. "
     "IP whitelisting (0.0.0.0/0) documented as prerequisite."),
    ("🔄", "CORS / Failed to Fetch", ACCENT2,
     "Problem", "Resolution",
     "Frontend calling http://localhost:5000/api directly caused 'Failed to fetch' "
     "errors when ports mismatched or backend was down.",
     "Added Vite server.proxy in vite.config.js to forward /api → localhost:5000. "
     "Frontend now uses relative /api path — no CORS issues, no hardcoded URLs."),
    ("📊", "Velocity Classification Edge Cases", PURPLE,
     "Problem", "Resolution",
     "salesVelocity = 0 caused division-by-zero in daysToSellOut formula. "
     "Also needed to handle items already expired (negative daysUntilExpiry).",
     "Guard: velocity=0 → daysToSellOut = Infinity (will never sell). "
     "Skip items where daysUntilExpiry ≤ 0 in the engine loop."),
]

for i, (icon, head, clr, lbl1, lbl2, prob, res) in enumerate(challenges):
    col = i % 2; row = i // 2
    lx = 0.45 + col * 6.45
    ty = 1.52 + row * 2.78
    rect(s, lx, ty, 6.1, 2.55, CARD, line_clr=BORDER, line_w=0.4)
    rect(s, lx, ty, 0.06, 2.55, clr)
    txt(s, icon + "  " + head, lx+0.2, ty+0.12, 5.7, 0.42, size=14, bold=True, color=clr)
    rect(s, lx+0.18, ty+0.62, 5.76, 0.02, BORDER)
    txt(s, "⚠ " + lbl1 + ":", lx+0.18, ty+0.72, 1.2, 0.3, size=11, bold=True, color=RED)
    txt(s, prob, lx+0.18, ty+0.72, 5.76, 0.78, size=12, color=OFFWHITE, wrap=True)
    txt(s, "✔ " + lbl2 + ":", lx+0.18, ty+1.55, 1.2, 0.3, size=11, bold=True, color=ACCENT)
    txt(s, res,  lx+0.18, ty+1.55, 5.76, 0.95, size=12, color=LGRAY, wrap=True)

footer(s, 11)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — FUTURE SCOPE
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s); section_header(s, "Future Scope", "Next-phase enhancements and product roadmap")

phases = [
    ("Phase 1\nNear-Term", ACCENT, [
        ("🤖", "Auto-Transfer Execution",   "System automatically executes approved low-risk transfers without manual confirmation"),
        ("📱", "Mobile Application",         "React Native app for store managers to approve transfers on the go"),
        ("🔔", "Push Notifications",         "Real-time SMS/email alerts when items enter the critical expiry window"),
    ]),
    ("Phase 2\nMid-Term", ACCENT2, [
        ("🧠", "ML Velocity Prediction",     "LSTM model to predict future sales velocity based on seasonal patterns"),
        ("📊", "Analytics Dashboard",        "Historical waste reduction reports, ROI metrics, and trend charts"),
        ("🔗", "ERP Integration",            "REST connectors for SAP, Tally, and popular POS systems"),
    ]),
    ("Phase 3\nLong-Term", GOLD, [
        ("🌐", "Multi-Chain Support",        "Manage thousands of stores across multiple grocery chains from one platform"),
        ("🚛", "Logistics Optimization",     "Route optimization for physical transfer trucks using Google Maps API"),
        ("♻️", "NGO / Food Bank Connect",   "Redirect unsellable items to food banks automatically before expiry"),
    ]),
]

for ci, (phase_lbl, clr, items) in enumerate(phases):
    lx = 0.45 + ci * 4.3
    rect(s, lx, 1.48, 4.1, 5.55, CARD, line_clr=clr, line_w=0.8)
    rect(s, lx, 1.48, 4.1, 0.07, clr)
    txt(s, phase_lbl, lx, 1.55, 4.1, 0.55,
        size=13, bold=True, color=clr, align=PP_ALIGN.CENTER)
    rect(s, lx+0.3, 2.18, 3.5, 0.03, BORDER)
    for ri, (icon, feat, desc) in enumerate(items):
        ty2 = 2.32 + ri * 1.55
        rect(s, lx+0.2, ty2, 3.72, 1.38, PANEL, line_clr=BORDER, line_w=0.3)
        txt(s, icon, lx+0.28, ty2+0.12, 0.6, 0.55, size=22)
        txt(s, feat, lx+0.92, ty2+0.1, 2.88, 0.42, size=13, bold=True, color=WHITE)
        txt(s, desc, lx+0.28, ty2+0.58, 3.5, 0.72, size=12, color=LGRAY, wrap=True)

footer(s, 12)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — CONCLUSION
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s); section_header(s, "Conclusion", "Summary of what we built, learned, and achieved")

# Summary quote
rect(s, 0.45, 1.48, 12.43, 0.88, RGBColor(0x06, 0x2A, 0x1F), line_clr=ACCENT, line_w=0.7)
txt(s, '"  FoodRedist demonstrates that a well-designed MERN stack application can solve a real-world supply chain '
    'problem — intelligently, efficiently, and at scale.  "',
    0.65, 1.55, 12.1, 0.72, size=14.5, color=RGBColor(0xA7, 0xF3, 0xD0), italic=True, align=PP_ALIGN.CENTER)

# Key takeaways
takeaways = [
    ("✅", "Problem Solved",        ACCENT,
     "Automated identification of surplus stock in slow stores and matching with fast-selling stores before expiry"),
    ("✅", "Technology Applied",     ACCENT2,
     "Full MERN stack with Mongoose models, Express REST API, Velocity Engine service, and React dashboard"),
    ("✅", "Algorithm Implemented",  GOLD,
     "Greedy velocity-based classification: daysToSellOut vs daysUntilExpiry to compute surplus and demand gap"),
    ("✅", "Data Integrity Ensured", PURPLE,
     "Atomic MongoDB operations for all transfer writes — no partial state; complete TransferLog audit trail"),
    ("✅", "Real Results Achieved",  RGBColor(0x34, 0xD3, 0x99),
     "7 transfer recommendations generated from seed data — 4 items, 3 stores, 0 projected waste"),
    ("✅", "Learnings Gained",       LGRAY,
     "Replica set transactions, Vite proxy config, velocity edge cases, CORS handling, and COM-based PDF export"),
]

for i, (icon, head, clr, desc) in enumerate(takeaways):
    col = i % 2; row = i // 2
    lx = 0.45 + col * 6.45
    ty = 2.55 + row * 1.52
    rect(s, lx, ty, 6.1, 1.35, CARD, line_clr=BORDER, line_w=0.4)
    rect(s, lx, ty, 0.06, 1.35, clr)
    txt(s, icon + "  " + head, lx+0.2, ty+0.12, 5.7, 0.38, size=14, bold=True, color=clr)
    txt(s, desc, lx+0.2, ty+0.55, 5.75, 0.72, size=12.5, color=OFFWHITE, wrap=True)

footer(s, 13)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Q&A / THANK YOU
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
rect(s, 0, 0, 13.33, 7.5, BG)
rect(s, 0, 0, 13.33, 0.07, ACCENT)
rect(s, 0, 7.43, 13.33, 0.07, ACCENT)

# Decorative circles
for r_size, clr, op in [(4.5, ACCENT, 0.03),(3.2, ACCENT2, 0.04),(1.8, ACCENT, 0.06)]:
    cx = (13.33 - r_size) / 2
    cy = (7.5 - r_size) / 2
    sh = s.shapes.add_shape(1, Inches(cx), Inches(cy), Inches(r_size), Inches(r_size))
    sh.fill.background()
    sh.line.color.rgb = clr
    sh.line.width = Pt(0.6)

txt(s, "?", 0, 2.1, 13.33, 2.5, size=130, bold=True,
    color=RGBColor(0x0F, 0x2A, 0x1F), align=PP_ALIGN.CENTER)

txt(s, "Questions & Answers", 0, 1.55, 13.33, 0.72,
    size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Open floor — we welcome your questions!", 0, 2.25, 13.33, 0.48,
    size=17, color=LGRAY, align=PP_ALIGN.CENTER, italic=True)

rect(s, 3.5, 2.88, 6.33, 0.04, ACCENT)

txt(s, "Thank You", 0, 3.05, 13.33, 0.72,
    size=42, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
txt(s, "for your time and attention", 0, 3.75, 13.33, 0.45,
    size=17, color=LGRAY, align=PP_ALIGN.CENTER)

# Team box
rect(s, 3.0, 4.42, 7.33, 1.82, PANEL, line_clr=ACCENT, line_w=0.9)
txt(s, TEAM_NAME, 3.0, 4.52, 7.33, 0.48,
    size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
rect(s, 4.0, 4.98, 5.33, 0.03, BORDER)
txt(s, "Saurav          Tushar          Ankit          Samyank",
    3.0, 5.05, 7.33, 0.45,
    size=16, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "B.Tech Computer Science  |  2025–26", 3.0, 5.52, 7.33, 0.35,
    size=12, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

txt(s, PROJECT + " — Perishable Food Stock Redistribution System",
    0, 6.35, 13.33, 0.38,
    size=12, color=GRAY, align=PP_ALIGN.CENTER)

rect(s, 0, 7.22, 13.33, 0.28, PANEL)
rect(s, 0, 7.22, 13.33, 0.025, ACCENT)
txt(s, TEAM_NAME + "  |  " + TEAM_MEMBERS, 0, 7.26, 13.33, 0.22,
    size=10, color=GRAY, align=PP_ALIGN.CENTER)


# ── Save PPTX ────────────────────────────────────────────────────────────────
pptx_path = r"c:\Users\sautr\food_app\FoodRedist_Corporate.pptx"
prs.save(pptx_path)
print(f"PPTX saved: {pptx_path}")
print(f"Slides: {len(prs.slides)}")
