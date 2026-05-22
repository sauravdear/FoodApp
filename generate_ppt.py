from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Colour palette ──────────────────────────────────────────────────────────
DARK_BG   = RGBColor(0x03, 0x07, 0x12)   # #030712
CARD_BG   = RGBColor(0x1F, 0x29, 0x37)   # #1F2937
GREEN     = RGBColor(0x10, 0xB9, 0x81)   # emerald-500
GREEN_DK  = RGBColor(0x05, 0x96, 0x69)   # emerald-600
RED       = RGBColor(0xEF, 0x44, 0x44)
ORANGE    = RGBColor(0xF9, 0x73, 0x16)
YELLOW    = RGBColor(0xEA, 0xB3, 0x08)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0x9C, 0xA3, 0xAF)
LIGHT     = RGBColor(0xF3, 0xF4, 0xF6)

blank_layout = prs.slide_layouts[6]   # completely blank


# ── Helper utilities ─────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_para(tf, text, font_size=16, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, space_before=6, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def full_bg(slide, color=DARK_BG):
    add_rect(slide, 0, 0, 13.33, 7.5, color)


def slide_header(slide, title, subtitle=None, accent=GREEN):
    # top accent bar
    add_rect(slide, 0, 0, 13.33, 0.08, accent)
    add_text(slide, title, 0.5, 0.2, 12, 0.7,
             font_size=36, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.5, 0.85, 12, 0.45,
                 font_size=17, color=GRAY, italic=True)
    # bottom accent line
    add_rect(slide, 0, 7.42, 13.33, 0.08, accent)


def icon_card(slide, icon, heading, body_lines, l, t, w=3.8, h=2.2,
              accent=GREEN):
    add_rect(slide, l, t, w, h, CARD_BG)
    add_rect(slide, l, t, w, 0.06, accent)          # top stripe
    add_text(slide, icon, l+0.1, t+0.1, 0.7, 0.5, font_size=26)
    add_text(slide, heading, l+0.85, t+0.12, w-1.0, 0.45,
             font_size=15, bold=True, color=WHITE)
    txb = slide.shapes.add_textbox(
        Inches(l+0.15), Inches(t+0.62), Inches(w-0.3), Inches(h-0.75))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in body_lines:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(13)
        run.font.color.rgb = GRAY


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title / Cover
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
full_bg(s)

# gradient-like left panel
add_rect(s, 0, 0, 5.5, 7.5, GREEN_DK)
add_rect(s, 5.3, 0, 0.4, 7.5, GREEN)

# big icon
add_text(s, "🥗", 1.2, 0.8, 3, 1.4, font_size=90, align=PP_ALIGN.CENTER)

add_text(s, "FoodRedist", 0.2, 2.4, 5.1, 1.0,
         font_size=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Perishable Food", 0.2, 3.3, 5.1, 0.55,
         font_size=22, color=LIGHT, align=PP_ALIGN.CENTER)
add_text(s, "Stock Redistribution System", 0.2, 3.82, 5.1, 0.55,
         font_size=22, color=LIGHT, align=PP_ALIGN.CENTER)

add_text(s, "MERN Stack Project", 0.2, 4.55, 5.1, 0.4,
         font_size=14, color=RGBColor(0xA7, 0xF3, 0xD0), align=PP_ALIGN.CENTER)

# right side — tagline
add_text(s, "Khaana Expire Hone Se Bachao,", 5.9, 1.5, 7.0, 0.65,
         font_size=28, bold=True, color=WHITE)
add_text(s, "Sahi Store Tak Pahuncho!", 5.9, 2.1, 7.0, 0.65,
         font_size=28, bold=True, color=GREEN)

add_text(s,
    "Ek aisa system jo track karta hai perishable\n"
    "food items ko, calculate karta hai har store ki\n"
    "selling speed, aur suggest karta hai ki kahan\n"
    "se kahan transfer karna chahiye — expire hone\n"
    "se pehle.",
    5.9, 2.95, 6.9, 2.5, font_size=16, color=GRAY)

# tags
for i, tag in enumerate(["MongoDB", "Express.js", "React", "Node.js"]):
    add_rect(s, 5.9 + i*1.65, 6.6, 1.5, 0.5, GREEN_DK)
    add_text(s, tag, 5.9 + i*1.65, 6.6, 1.5, 0.5,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
full_bg(s)
slide_header(s, "Problem Kya Hai?", "Grocery stores mein roz ka struggle — waste vs. shortage", accent=RED)

add_text(s, "😟 Real-World Problem", 0.5, 1.4, 12, 0.5,
         font_size=22, bold=True, color=RED)

problems = [
    ("🗑️", "Ek store mein 100 units milk padi hai jo kal expire hogi",
           "Par wahan customers kum aate hain — sab waste ho jaayega"),
    ("🏃", "Doosre store mein milk ki demand zyada hai",
           "Par stock sirf 5 units bachi hai — customers khali haath jaate hain"),
    ("📊", "Koi system nahi hai jo ye bataye",
           "Manager ko manually pata karna padta hai — tab tak der ho jaati hai"),
    ("💸", "Result: Crores ka food waste har saal",
           "Aur saath mein unhappy customers bhi — double loss!"),
]

for i, (icon, line1, line2) in enumerate(problems):
    col = i % 2
    row = i // 2
    lx = 0.5 + col * 6.5
    ty = 2.0 + row * 2.3
    add_rect(s, lx, ty, 6.1, 2.05, CARD_BG)
    add_rect(s, lx, ty, 0.06, 2.05, RED)
    add_text(s, icon, lx+0.2, ty+0.3, 0.7, 0.7, font_size=28)
    add_text(s, line1, lx+0.95, ty+0.2, 5.0, 0.55, font_size=14, bold=True, color=WHITE)
    add_text(s, line2, lx+0.95, ty+0.75, 5.0, 0.55, font_size=13, color=GRAY)

add_rect(s, 0.5, 6.6, 12.3, 0.55, RGBColor(0x45, 0x10, 0x10))
add_text(s,
    "✅  FoodRedist iska solution hai — real-time tracking + automated recommendations",
    0.7, 6.63, 12.0, 0.45, font_size=15, bold=True, color=RGBColor(0xFC, 0xA5, 0xA5))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Solution Overview
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
full_bg(s)
slide_header(s, "Hamaara Solution", "FoodRedist — Smart redistribution, zero waste", accent=GREEN)

add_text(s, "System kaise kaam karta hai — 3 simple steps:", 0.5, 1.3, 12, 0.45,
         font_size=18, color=GRAY, italic=True)

steps = [
    ("1️⃣", "TRACK", "Real-time monitoring",
     ["Har store ki inventory track hoti hai", "Expiry dates monitor hoti hain", "Sales velocity calculate hoti hai"]),
    ("2️⃣", "ANALYZE", "Velocity Engine",
     ["Algorithm decide karta hai — kaunsa store slow hai", "Surplus units calculate hote hain", "Best match dhundha jaata hai"]),
    ("3️⃣", "ACT", "1-Click Transfer",
     ["Manager ko recommendation milti hai", "Confirm karo — transfer ho jaata hai", "Audit log automatically banta hai"]),
]

for i, (num, heading, sub, bullets) in enumerate(steps):
    lx = 0.5 + i * 4.25
    add_rect(s, lx, 1.9, 4.0, 4.9, CARD_BG)
    add_rect(s, lx, 1.9, 4.0, 0.08, GREEN)
    add_text(s, num, lx+0.15, 2.05, 0.8, 0.7, font_size=34)
    add_text(s, heading, lx+1.0, 2.05, 2.8, 0.5, font_size=20, bold=True, color=GREEN)
    add_text(s, sub, lx+0.2, 2.62, 3.6, 0.4, font_size=13, color=GRAY, italic=True)
    add_rect(s, lx+0.2, 3.1, 3.6, 0.02, GRAY)
    txb = slide.shapes.add_textbox(Inches(lx+0.2), Inches(3.22), Inches(3.6), Inches(3.3)) if False else \
          s.shapes.add_textbox(Inches(lx+0.2), Inches(3.22), Inches(3.6), Inches(3.3))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for j, b in enumerate(bullets):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = "  ▸  " + b
        run.font.size = Pt(14)
        run.font.color.rgb = LIGHT

    # connector arrow (except last)
    if i < 2:
        add_text(s, "→", lx+4.0, 3.9, 0.5, 0.5, font_size=28, color=GREEN, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Tech Stack (MERN)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
full_bg(s)
slide_header(s, "Tech Stack — MERN", "Kaunsi technologies use hui hain aur kyun?")

techs = [
    ("🍃", "MongoDB", "Database", GREEN,
     ["NoSQL document database",
      "Flexible schema — perfect for inventory",
      "Mongoose ORM se models banaye",
      "Atlas cloud ya local dono pe run karta hai"]),
    ("⚡", "Express.js", "Backend Framework", ORANGE,
     ["Node.js ke upar REST API banata hai",
      "Routes: /api/inventory, /api/stores",
      "Middleware: CORS, JSON parsing",
      "Clean MVC architecture follow kiya"]),
    ("⚛️", "React.js", "Frontend UI", RGBColor(0x61, 0xDB, 0xFB),
     ["Component-based UI",
      "Context API for global state",
      "React Router for navigation",
      "Tailwind CSS for dark theme styling"]),
    ("🟩", "Node.js", "Runtime", RGBColor(0x68, 0xA0, 0x63),
     ["JavaScript runtime on server",
      "npm packages manage karta hai",
      "Async/await for DB operations",
      "ES Modules (import/export) syntax"]),
]

for i, (icon, name, role, clr, points) in enumerate(techs):
    col = i % 2
    row = i // 2
    lx = 0.4 + col * 6.45
    ty = 1.45 + row * 2.85
    add_rect(s, lx, ty, 6.1, 2.6, CARD_BG)
    add_rect(s, lx, ty, 6.1, 0.07, clr)
    add_text(s, icon, lx+0.15, ty+0.18, 0.75, 0.65, font_size=32)
    add_text(s, name, lx+0.95, ty+0.15, 3.5, 0.45, font_size=20, bold=True, color=clr)
    add_text(s, role, lx+0.95, ty+0.55, 3.5, 0.35, font_size=13, color=GRAY, italic=True)
    txb = s.shapes.add_textbox(Inches(lx+0.2), Inches(ty+0.98), Inches(5.7), Inches(1.55))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    for j, pt in enumerate(points):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run(); run.text = "✓  " + pt
        run.font.size = Pt(13); run.font.color.rgb = GRAY


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Database Schema
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
full_bg(s)
slide_header(s, "Database Schema", "3 Mongoose models — data kaise store hota hai")

schemas = [
    ("📦", "FoodItem", GREEN,
     ["name        → item ka naam (Greek Yogurt)",
      "sku         → unique product code",
      "batchNumber → manufacturing batch",
      "expirationDate → expire kab hoga",
      "basePrice   → kitne ka hai"]),
    ("🏪", "Store", ORANGE,
     ["storeName   → store ka naam",
      "location    → address, city, state",
      "inventory[] → array of items:",
      "  • foodItemId  (ref → FoodItem)",
      "  • currentStock (number)",
      "  • salesVelocity (units/day)"]),
    ("📋", "TransferLog", RGBColor(0xA7, 0x8B, 0xFA),
     ["foodItemId       → kaunsa item",
      "sourceStoreId    → kahan se",
      "destinationStoreId → kahan tak",
      "quantity         → kitna",
      "status           → Pending/Completed",
      "timestamp        → kab hua"]),
]

for i, (icon, name, clr, fields) in enumerate(schemas):
    lx = 0.4 + i * 4.25
    add_rect(s, lx, 1.45, 4.0, 5.5, CARD_BG)
    add_rect(s, lx, 1.45, 4.0, 0.07, clr)
    add_text(s, icon + "  " + name, lx+0.2, 1.6, 3.6, 0.55,
             font_size=18, bold=True, color=clr)
    add_rect(s, lx+0.2, 2.28, 3.6, 0.03, CARD_BG)

    txb = s.shapes.add_textbox(Inches(lx+0.2), Inches(2.35), Inches(3.65), Inches(4.45))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    for j, field in enumerate(fields):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_before = Pt(7)
        run = p.add_run()
        run.text = field
        run.font.size = Pt(13)
        run.font.name = "Consolas"
        run.font.color.rgb = GRAY if field.startswith("  ") else LIGHT

    if i < 2:
        add_text(s, "🔗", lx+4.0, 3.85, 0.5, 0.5, font_size=22, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Velocity Engine
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
full_bg(s)
slide_header(s, "Velocity Engine — Dimaag of the System",
             "services/velocityEngine.js — yahi decide karta hai kahan se kahan jaaye")

# Formula box
add_rect(s, 0.5, 1.4, 12.3, 1.1, RGBColor(0x06, 0x4E, 0x3B))
add_rect(s, 0.5, 1.4, 0.07, 1.1, GREEN)
add_text(s, "🧮  Core Formula:", 0.75, 1.48, 4, 0.4, font_size=15, bold=True, color=GREEN)
add_text(s, "daysToSellOut  =  currentStock  ÷  salesVelocity",
         4.5, 1.48, 8.5, 0.4, font_size=20, bold=True, color=WHITE)
add_text(s, "Agar  daysToSellOut  >  daysUntilExpiry  →  SLOW store  |  Velocity ≥ threshold AND stock low  →  FAST store",
         0.75, 1.88, 12.0, 0.4, font_size=13, color=RGBColor(0xA7, 0xF3, 0xD0))

# SLOW vs FAST cards
for clr, title, icon, items, lx in [
    (RED, "SLOW Store", "🐢", [
        "salesVelocity kum hai (< threshold)",
        "daysToSellOut > daysUntilExpiry",
        "Matlab: ye store expire hone se pehle",
        "  stock khatam nahi kar paayega",
        "Surplus units calculate hote hain:",
        "  surplus = stock − (velocity × daysLeft)",
    ], 0.5),
    (GREEN, "FAST Store", "🚀", [
        "salesVelocity >= threshold",
        "currentStock <= lowStockGuard (20 units)",
        "Matlab: yahan demand zyada hai",
        "  par stock khatam ho rahi hai",
        "Demand gap calculate hota hai:",
        "  gap = (velocity × daysLeft) − stock",
    ], 6.9),
]:
    add_rect(s, lx, 2.7, 5.9, 4.4, CARD_BG)
    add_rect(s, lx, 2.7, 5.9, 0.08, clr)
    add_text(s, icon + "  " + title, lx+0.2, 2.85, 5.0, 0.55,
             font_size=19, bold=True, color=clr)
    txb = s.shapes.add_textbox(Inches(lx+0.2), Inches(3.5), Inches(5.55), Inches(3.4))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    for j, item in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(14)
        run.font.name = "Consolas" if item.startswith("  ") else "Calibri"
        run.font.color.rgb = GRAY if item.startswith("  ") else LIGHT

add_text(s, "⟶", 6.45, 4.5, 0.7, 0.7, font_size=36, color=YELLOW, align=PP_ALIGN.CENTER)
add_text(s, "Match!", 6.3, 5.1, 1.0, 0.4, font_size=13, bold=True,
         color=YELLOW, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — API Endpoints
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
full_bg(s)
slide_header(s, "API Endpoints", "Express routes — frontend in se data maangti hai")

endpoints = [
    ("GET", "/api/inventory/recommendations", GREEN,
     "Velocity engine run karta hai",
     ["Query params: expiryWindowDays, velocityThreshold, lowStockGuard",
      "Returns: expiring items + recommended transfers list",
      "Dashboard pe Stock Alerts yahi se aate hain"]),
    ("POST", "/api/inventory/transfer", ORANGE,
     "Transfer execute karta hai",
     ["Body: { foodItemId, sourceStoreId, destinationStoreId, quantity }",
      "Step 1: Source store ka stock decrement hota hai",
      "Step 2: Destination store ka stock increment hota hai",
      "Step 3: TransferLog record create hota hai"]),
    ("GET", "/api/inventory/transfers", RGBColor(0xA7, 0x8B, 0xFA),
     "Puri history dikhata hai",
     ["Populated data: item name, store names",
      "Sorted by timestamp (newest first)",
      "Transfers page pe yahi dikhta hai"]),
    ("GET", "/api/stores", RGBColor(0x61, 0xDB, 0xFB),
     "Sabhi stores list karta hai",
     ["Har store ki full inventory",
      "FoodItem details bhi populate hoti hain",
      "Used for store management"]),
    ("PATCH", "/api/stores/:id/inventory/:fid/velocity", YELLOW,
     "Velocity update karta hai",
     ["Ek specific item ki salesVelocity update hoti hai",
      "Velocity engine ke baad call hota hai",
      "Daily sales data se calculate karke set karo"]),
    ("GET", "/api/health", GRAY,
     "Server status check",
     ["Returns: { status: 'ok', timestamp }",
      "Debugging ke liye useful",
      "Deployment monitoring mein use hota hai"]),
]

for i, (method, path, clr, desc, details) in enumerate(endpoints):
    col = i % 2
    row = i // 2
    lx = 0.4 + col * 6.45
    ty = 1.45 + row * 1.95
    add_rect(s, lx, ty, 6.1, 1.8, CARD_BG)
    add_rect(s, lx, ty, 0.07, 1.8, clr)

    add_rect(s, lx+0.2, ty+0.18, 0.85, 0.38, clr)
    add_text(s, method, lx+0.2, ty+0.18, 0.85, 0.38,
             font_size=11, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_text(s, path, lx+1.15, ty+0.15, 4.7, 0.42,
             font_size=13, bold=True, color=WHITE)
    add_text(s, desc, lx+0.2, ty+0.6, 5.7, 0.38,
             font_size=12, color=GRAY, italic=True)

    detail_text = "  •  " + ("  •  ".join(details[:2]))
    add_text(s, detail_text, lx+0.2, ty+0.98, 5.7, 0.62,
             font_size=11, color=RGBColor(0x6B, 0x72, 0x80))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Frontend UI Walkthrough
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
full_bg(s)
slide_header(s, "Frontend UI — Kya Dikhta Hai?",
             "React components — user kya dekhta aur karta hai")

components = [
    ("🏠", "Navbar", GREEN,
     "Top navigation bar",
     ["Dashboard aur Transfers ke beech navigate karo",
      "Active page highlight hoti hai",
      "FoodRedist logo dono pe dikhta hai"]),
    ("📊", "StockAlerts", RED,
     "Dashboard — Left Panel",
     ["Har expiring item ka card dikhta hai",
      "Color: Red (1 din), Orange (2 din), Yellow (3 din)",
      "Item name, SKU, batch, velocity info shown hoti hai",
      "Recommended transfer quantity dikhti hai"]),
    ("✅", "TransferPanel", ORANGE,
     "Dashboard — Right Panel",
     ["Kisi bhi alert pe click karo → panel mein details aate hain",
      "Quantity edit kar sakte ho",
      "Confirm Transfer → API call → DB update",
      "Success/error message instantly dikhta hai"]),
    ("🔄", "Transfers Page", RGBColor(0xA7, 0x78, 0xFA),
     "History Page",
     ["Saare completed transfers ki list",
      "Har entry mein: item, from store, to store, qty, status",
      "Status badge: green (Completed), yellow (Pending)",
      "Newest pehle dikhta hai"]),
]

for i, (icon, name, clr, role, pts) in enumerate(components):
    col = i % 2
    row = i // 2
    lx = 0.4 + col * 6.45
    ty = 1.45 + row * 2.8
    add_rect(s, lx, ty, 6.1, 2.55, CARD_BG)
    add_rect(s, lx, ty, 6.1, 0.07, clr)
    add_text(s, icon, lx+0.15, ty+0.18, 0.7, 0.6, font_size=28)
    add_text(s, name, lx+0.9, ty+0.15, 3.5, 0.45, font_size=18, bold=True, color=clr)
    add_text(s, role, lx+0.9, ty+0.55, 4.8, 0.38, font_size=13, color=GRAY, italic=True)
    txb = s.shapes.add_textbox(Inches(lx+0.2), Inches(ty+1.0), Inches(5.75), Inches(1.45))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    for j, pt in enumerate(pts):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run(); run.text = "▸  " + pt
        run.font.size = Pt(13); run.font.color.rgb = LIGHT


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Data Flow (End to End)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
full_bg(s)
slide_header(s, "End-to-End Data Flow",
             "Ek transfer request ka poora safar — DB se UI tak")

flow_steps = [
    ("🗄️", "MongoDB Atlas /\nLocal MongoDB", "Database", GRAY),
    ("⚙️", "Velocity Engine\nvelocityEngine.js", "Backend Service", GREEN),
    ("🚦", "Express\nController", "inventoryController.js", ORANGE),
    ("🌐", "REST API\n/api/inventory", "HTTP Layer", RGBColor(0xA7, 0x8B, 0xFA)),
    ("⚛️", "React\nAppContext", "State Management", RGBColor(0x61, 0xDB, 0xFB)),
    ("🖥️", "Dashboard\nUI", "User Interface", YELLOW),
]

box_w = 1.85
gap   = 0.25
start = 0.4

for i, (icon, label, sublabel, clr) in enumerate(flow_steps):
    lx = start + i * (box_w + gap)
    add_rect(s, lx, 2.2, box_w, 2.8, CARD_BG)
    add_rect(s, lx, 2.2, box_w, 0.07, clr)
    add_text(s, icon, lx + box_w/2 - 0.25, 2.35, 0.6, 0.6, font_size=28, align=PP_ALIGN.CENTER)
    add_text(s, label, lx+0.1, 2.98, box_w-0.2, 0.75,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, sublabel, lx+0.1, 3.75, box_w-0.2, 0.5,
             font_size=11, color=clr, align=PP_ALIGN.CENTER, italic=True)
    if i < len(flow_steps)-1:
        add_text(s, "→", lx+box_w+0.02, 3.1, gap+0.2, 0.5,
                 font_size=22, color=GRAY, align=PP_ALIGN.CENTER)

# Step-by-step flow description
flow_desc = [
    "1️⃣  Manager dashboard open karta hai → React frontend /api/recommendations pe GET request bhejta hai",
    "2️⃣  Express controller velocityEngine.analyzeInventory() call karta hai",
    "3️⃣  Velocity Engine MongoDB se expiring items aur stores fetch karta hai, slow/fast classify karta hai",
    "4️⃣  Recommendations JSON mein wapas aate hain → AppContext mein store hote hain → StockAlerts render hota hai",
    "5️⃣  Manager ek alert select karta hai, quantity confirm karta hai → POST /api/inventory/transfer",
    "6️⃣  Controller source stock decrement, destination increment, TransferLog create karta hai → 201 response",
]

txb = s.shapes.add_textbox(Inches(0.5), Inches(5.25), Inches(12.3), Inches(2.0))
txb.word_wrap = True
tf = txb.text_frame; tf.word_wrap = True
for j, line in enumerate(flow_desc):
    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
    p.space_before = Pt(4)
    run = p.add_run(); run.text = line
    run.font.size = Pt(12.5)
    run.font.color.rgb = GRAY


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Live Demo / How to Run
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
full_bg(s)
slide_header(s, "Project Kaise Run Karte Hain?",
             "Setup se lekar live demo tak — step by step")

# Left: setup steps
add_rect(s, 0.4, 1.4, 6.1, 5.7, CARD_BG)
add_rect(s, 0.4, 1.4, 0.07, 5.7, GREEN)
add_text(s, "🚀  Quick Start", 0.65, 1.5, 5.5, 0.5,
         font_size=18, bold=True, color=GREEN)

setup_lines = [
    ("OPTION A — One Click", True, GREEN),
    ("Double-click  start.bat  in food_app/", False, LIGHT),
    ("MongoDB + Backend + Frontend sab start!", False, GRAY),
    ("", False, GRAY),
    ("OPTION B — Manual", True, ORANGE),
    ("# Terminal 1 — Backend", False, GRAY),
    ("cd backend  &&  node server.js", False, RGBColor(0x86, 0xEF, 0xAC)),
    ("", False, GRAY),
    ("# Terminal 2 — Frontend", False, GRAY),
    ("cd frontend  &&  npm run dev", False, RGBColor(0x86, 0xEF, 0xAC)),
    ("", False, GRAY),
    ("URLs:", True, LIGHT),
    ("Frontend  →  http://localhost:5173", False, RGBColor(0x61, 0xDB, 0xFB)),
    ("Backend   →  http://localhost:5000", False, ORANGE),
    ("Health    →  /api/health", False, GRAY),
]

txb = s.shapes.add_textbox(Inches(0.65), Inches(2.05), Inches(5.7), Inches(4.9))
txb.word_wrap = True
tf = txb.text_frame; tf.word_wrap = True
for j, (text, bold, clr) in enumerate(setup_lines):
    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
    p.space_before = Pt(5)
    run = p.add_run(); run.text = text
    run.font.size = Pt(13.5)
    run.font.bold = bold
    run.font.name = "Consolas" if text.startswith("cd ") or text.startswith("#") or text.startswith("http") else "Calibri"
    run.font.color.rgb = clr

# Right: demo steps
add_rect(s, 6.8, 1.4, 6.1, 5.7, CARD_BG)
add_rect(s, 6.8, 1.4, 0.07, 5.7, YELLOW)
add_text(s, "🎯  Demo Steps", 7.05, 1.5, 5.5, 0.5,
         font_size=18, bold=True, color=YELLOW)

demo_steps = [
    ("1", "Dashboard kholo", "localhost:5173 browser mein open karo", GREEN),
    ("2", "Stock Alerts dekho", "7 alerts dikhenge — red/orange/yellow urgency se", RED),
    ("3", "Alert select karo", "Kisi bhi card pe click karo — Transfer Panel khulega", ORANGE),
    ("4", "Quantity confirm karo", "Default recommended qty hai, change bhi kar sakte ho", YELLOW),
    ("5", "Confirm Transfer dabao", "Success message aayega — stock update ho jaayega", GREEN),
    ("6", "Transfers page dekho", "Saara history dikhega completed transfers ka", RGBColor(0xA7, 0x8B, 0xFA)),
]

for j, (num, heading, detail, clr) in enumerate(demo_steps):
    ty = 2.08 + j * 0.83
    add_rect(s, 7.05, ty, 0.42, 0.42, clr)
    add_text(s, num, 7.05, ty, 0.42, 0.42, font_size=14, bold=True,
             color=DARK_BG, align=PP_ALIGN.CENTER)
    add_text(s, heading, 7.58, ty+0.01, 5.0, 0.28, font_size=13, bold=True, color=WHITE)
    add_text(s, detail,  7.58, ty+0.29, 5.0, 0.32, font_size=12, color=GRAY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Why We Built This / Impact
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
full_bg(s)
slide_header(s, "Kyun Banaya? — Impact & Vision",
             "Ek small project, ek badi problem ka solution")

add_text(s,
    "India mein har saal ₹92,000 crore ka food waste hota hai.\n"
    "FoodRedist iska ek hissa rokne ki koshish karta hai.",
    0.5, 1.35, 12.3, 0.8,
    font_size=17, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

impacts = [
    ("📉", "Food Waste Kum", GREEN,
     "Expiring items ko sahi store pe\nbhejke waste eliminate karo"),
    ("😊", "Customer Satisfaction", ORANGE,
     "Fast stores mein stock kabhi\nkhatam nahi hoga"),
    ("💰", "Revenue Protect", YELLOW,
     "Jo item waste hota, usse bhi\nbech ke revenue save karo"),
    ("📊", "Data-Driven Decisions", RGBColor(0xA7, 0x8B, 0xFA),
     "Gut feeling nahi — algorithm\nbatata hai kya karna hai"),
    ("⚡", "Real-Time Action", RGBColor(0x61, 0xDB, 0xFB),
     "Alert aaya, click kiya, transfer\nhogaya — minutes mein"),
    ("🔮", "Future Scope", GRAY,
     "Auto-transfer, ML predictions,\nmobile app, multi-chain support"),
]

for i, (icon, heading, clr, detail) in enumerate(impacts):
    col = i % 3
    row = i // 2
    lx = 0.4 + col * 4.25
    ty = 2.35 + row * 2.4
    add_rect(s, lx, ty, 4.0, 2.1, CARD_BG)
    add_rect(s, lx, ty, 4.0, 0.07, clr)
    add_text(s, icon, lx+0.2, ty+0.2, 0.65, 0.65, font_size=30)
    add_text(s, heading, lx+0.9, ty+0.25, 2.9, 0.45, font_size=15, bold=True, color=clr)
    add_text(s, detail, lx+0.2, ty+0.82, 3.65, 1.1, font_size=13, color=GRAY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Thank You / Summary
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
full_bg(s)
add_rect(s, 0, 0, 13.33, 7.5, DARK_BG)
add_rect(s, 0, 0, 13.33, 0.1, GREEN)
add_rect(s, 0, 7.4, 13.33, 0.1, GREEN)

add_text(s, "🙏", 0, 1.2, 13.33, 1.2, font_size=64, align=PP_ALIGN.CENTER)
add_text(s, "Shukriya!", 0, 2.35, 13.33, 0.9,
         font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "FoodRedist — Perishable Food Stock Redistribution System",
         0, 3.2, 13.33, 0.55,
         font_size=20, color=GREEN, align=PP_ALIGN.CENTER)

summary = "MongoDB + Express + React + Node.js  |  Velocity Engine  |  Real-time Alerts  |  Atomic Transfers"
add_text(s, summary, 0, 3.88, 13.33, 0.45,
         font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

add_rect(s, 2.5, 4.55, 8.33, 0.03, CARD_BG)

key_points = [
    "✅ Problem solve kiya: food waste + stock shortage",
    "✅ 3 MongoDB models, REST API, Velocity Engine",
    "✅ React dashboard with real-time alerts",
    "✅ 1-click transfer with audit log",
]
txb = s.shapes.add_textbox(Inches(2.5), Inches(4.8), Inches(8.33), Inches(2.0))
txb.word_wrap = True
tf = txb.text_frame; tf.word_wrap = True
for j, kp in enumerate(key_points):
    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(6)
    run = p.add_run(); run.text = kp
    run.font.size = Pt(15); run.font.color.rgb = GRAY


# ── Save ────────────────────────────────────────────────────────────────────
out = r"c:\Users\sautr\food_app\FoodRedist_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
